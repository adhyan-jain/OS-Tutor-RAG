"""Extract text and metadata from PowerPoint (.ppt/.pptx) files into Documents."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from src.schemas import Document


def extract_ppt(file_path: Path) -> Document:
    """Extract a Document from a single PPT/PPTX file.

    Each slide's title, bullet text, and speaker notes are captured as
    separate fields in a per-slide dict under ``metadata["slides"]``. The
    Document's ``text`` is the concatenation of slide titles and bullets
    (notes are kept in metadata only, since they are speaker-facing rather
    than slide content).

    Args:
        file_path: Path to the .ppt or .pptx file.

    Returns:
        A Document containing the concatenated slide text and slide-level metadata.
    """
    presentation = Presentation(file_path)

    slides_metadata: list[dict[str, object]] = []
    text_parts: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else None

        bullets: list[str] = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = "".join(run.text for run in paragraph.runs).strip()
                if paragraph_text:
                    bullets.append(paragraph_text)

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides_metadata.append(
            {
                "slide_number": slide_number,
                "title": title,
                "bullets": bullets,
                "notes": notes,
            }
        )

        if title:
            text_parts.append(title)
        text_parts.extend(bullets)

    return Document(
        doc_id=file_path.stem,
        source_path=str(file_path),
        text="\n".join(text_parts),
        metadata={
            "source_type": "pptx",
            "slide_count": len(presentation.slides),
            "slides": slides_metadata,
        },
    )
